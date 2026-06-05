"""内置工具：文档读取（Office / PDF 文本提取）

file_read 只能按文本读取，docx/xlsx/pdf 等二进制文档由本工具负责：
  - .docx        → python-docx：段落 + 表格（表格渲染为 Markdown）
  - .xlsx/.xlsm  → openpyxl(read_only)：sheet 列表 + 行数据（制表符分隔）
  - .xls         → xlrd：同 xlsx 形态
  - .pdf         → pypdf：按页提取文本，页边界标记
  - .pptx        → python-pptx：逐页文本
  - .doc/.ppt/.wps → 不支持二进制老格式，返回转换指引（另存为 docx/pptx）

输出 JSON 风格与 file_tool 一致：{"success", "format", "content", ...}，
content 超长截断并给出继续读取的提示（pdf 用 page_start/page_end，
xlsx 用 sheet 参数）。解析库为核心硬依赖；import 失败时防御性返回安装提示。
"""
from __future__ import annotations

import json
import os
from typing import Optional

from milu.tools.decorator import tool

# 单次返回内容上限（与 file_tool 全文读取一致）
_MAX_CONTENT_CHARS = 50000
# xlsx/xls 单 sheet 最大读取行数
_MAX_SHEET_ROWS = 500
# pdf 默认最多读取页数（未指定范围时）
_DEFAULT_PDF_PAGES = 20

# 老格式 → 转换指引
_LEGACY_HINTS = {
    ".doc": "请用 Word/WPS 打开后「另存为 .docx」，或用 LibreOffice 转换："
            "soffice --headless --convert-to docx 文件.doc",
    ".ppt": "请用 PowerPoint/WPS 打开后「另存为 .pptx」，或用 LibreOffice 转换："
            "soffice --headless --convert-to pptx 文件.ppt",
    ".wps": "请用 WPS 打开后「另存为 .docx」再读取",
}


def _err(msg: str, **extra) -> dict:
    return {"success": False, "error": msg, **extra}


def _truncate(text: str) -> tuple[str, bool]:
    """内容截断：返回 (内容, 是否截断)。"""
    if len(text) > _MAX_CONTENT_CHARS:
        return text[:_MAX_CONTENT_CHARS], True
    return text, False


# ── 各格式提取实现 ─────────────────────────────────────────


def _read_docx(path: str, **_) -> dict:
    try:
        from docx import Document
        from docx.document import Document as _DocxDocument
        from docx.table import Table
        from docx.text.paragraph import Paragraph
    except ImportError:
        return _err("缺少 python-docx，请安装：pip install python-docx")

    doc: "_DocxDocument" = Document(path)

    # 按文档顺序遍历段落与表格（body 子元素级遍历）
    parts: list[str] = []
    para_count = 0
    table_count = 0
    for element in doc.element.body:
        if element.tag.endswith("}p"):
            text = Paragraph(element, doc).text.strip()
            if text:
                parts.append(text)
                para_count += 1
        elif element.tag.endswith("}tbl"):
            table = Table(element, doc)
            table_count += 1
            rows = []
            for i, row in enumerate(table.rows):
                cells = [c.text.strip().replace("\n", " ") for c in row.cells]
                rows.append("| " + " | ".join(cells) + " |")
                if i == 0:  # Markdown 表头分隔线
                    rows.append("|" + "---|" * len(cells))
            parts.append("\n".join(rows))

    content, truncated = _truncate("\n\n".join(parts))
    result = {
        "success": True,
        "format": "docx",
        "paragraphs": para_count,
        "tables": table_count,
        "content": content,
    }
    if truncated:
        result["truncated"] = True
        result["hint"] = "文档较长已截断，如需更多内容可转存为文本后用 file_read 分段读取"
    return result


def _rows_to_text(rows_iter, max_rows: int) -> tuple[str, int, bool]:
    """通用：行迭代器 → 制表符分隔文本。返回 (文本, 实际行数, 是否截断)。"""
    lines = []
    truncated = False
    count = 0
    for row in rows_iter:
        if count >= max_rows:
            truncated = True
            break
        cells = ["" if v is None else str(v) for v in row]
        lines.append("\t".join(cells).rstrip())
        count += 1
    return "\n".join(lines), count, truncated


def _read_xlsx(path: str, sheet: str | None = None, **_) -> dict:
    try:
        from openpyxl import load_workbook
    except ImportError:
        return _err("缺少 openpyxl，请安装：pip install openpyxl")

    wb = load_workbook(path, read_only=True, data_only=True)
    try:
        sheet_names = wb.sheetnames
        if sheet is not None and sheet not in sheet_names:
            return _err(f"工作表不存在: {sheet}", sheets=sheet_names)

        target = sheet or sheet_names[0]
        ws = wb[target]
        text, row_count, truncated = _rows_to_text(
            ws.iter_rows(values_only=True), _MAX_SHEET_ROWS
        )
        content, char_truncated = _truncate(text)

        result = {
            "success": True,
            "format": "xlsx",
            "sheets": sheet_names,
            "sheet": target,
            "rows_read": row_count,
            "content": content,
        }
        if truncated or char_truncated:
            result["truncated"] = True
            result["hint"] = (
                f"工作表超过 {_MAX_SHEET_ROWS} 行或内容过长已截断；"
                "其他工作表用 sheet 参数指定读取"
            )
        elif len(sheet_names) > 1 and sheet is None:
            result["hint"] = "工作簿含多个工作表，默认读取第一个；其余用 sheet 参数指定"
        return result
    finally:
        wb.close()


def _read_xls(path: str, sheet: str | None = None, **_) -> dict:
    try:
        import xlrd
    except ImportError:
        return _err("缺少 xlrd，请安装：pip install xlrd")

    wb = xlrd.open_workbook(path)
    sheet_names = wb.sheet_names()
    if sheet is not None and sheet not in sheet_names:
        return _err(f"工作表不存在: {sheet}", sheets=sheet_names)

    target = sheet or sheet_names[0]
    ws = wb.sheet_by_name(target)
    rows_iter = (ws.row_values(i) for i in range(ws.nrows))
    text, row_count, truncated = _rows_to_text(rows_iter, _MAX_SHEET_ROWS)
    content, char_truncated = _truncate(text)

    result = {
        "success": True,
        "format": "xls",
        "sheets": sheet_names,
        "sheet": target,
        "rows_read": row_count,
        "content": content,
    }
    if truncated or char_truncated:
        result["truncated"] = True
        result["hint"] = f"工作表超过 {_MAX_SHEET_ROWS} 行或内容过长已截断"
    return result


def _read_pdf(path: str, page_start: int | None = None,
              page_end: int | None = None, **_) -> dict:
    try:
        from pypdf import PdfReader
    except ImportError:
        return _err("缺少 pypdf，请安装：pip install pypdf")

    reader = PdfReader(path)
    total = len(reader.pages)

    # 页码范围（1-based），默认读前 _DEFAULT_PDF_PAGES 页
    start = max(1, page_start or 1)
    default_end = min(total, start + _DEFAULT_PDF_PAGES - 1)
    end = min(total, page_end or default_end)
    if start > total:
        return _err(f"页码越界：起始页 {start}，文档共 {total} 页")

    parts = []
    for i in range(start - 1, end):
        text = (reader.pages[i].extract_text() or "").strip()
        parts.append(f"--- 第 {i + 1} 页 ---\n{text}")

    content, truncated = _truncate("\n\n".join(parts))
    result = {
        "success": True,
        "format": "pdf",
        "total_pages": total,
        "page_start": start,
        "page_end": end,
        "content": content,
    }
    if truncated or end < total:
        result["truncated"] = True
        result["hint"] = (
            f"文档共 {total} 页，本次返回第 {start}-{end} 页；"
            "用 page_start/page_end 参数继续读取后续页"
        )
    return result


def _read_pptx(path: str, **_) -> dict:
    try:
        from pptx import Presentation
    except ImportError:
        return _err("缺少 python-pptx，请安装：pip install python-pptx")

    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    texts.append(t)
        parts.append(f"--- 幻灯片 {i} ---\n" + "\n".join(texts))

    content, truncated = _truncate("\n\n".join(parts))
    result = {
        "success": True,
        "format": "pptx",
        "slides": len(parts),
        "content": content,
    }
    if truncated:
        result["truncated"] = True
        result["hint"] = "演示文稿较长已截断"
    return result


_READERS = {
    ".docx": _read_docx,
    ".xlsx": _read_xlsx,
    ".xlsm": _read_xlsx,
    ".xls": _read_xls,
    ".pdf": _read_pdf,
    ".pptx": _read_pptx,
}

DOC_EXTENSIONS = frozenset(_READERS) | frozenset(_LEGACY_HINTS)


# ── 工具定义 ──────────────────────────────────────────────


@tool(name="doc_read", description=(
    "读取 Office/PDF 文档并提取文本内容。支持格式：\n"
    "docx（段落+表格）、xlsx/xlsm/xls（sheet 参数指定工作表）、"
    "pdf（page_start/page_end 指定页码范围，默认前 20 页）、pptx（逐页文本）。\n"
    "读取 Word/Excel/PPT/PDF 等二进制文档必须用本工具，不要用 file_read。"
    ".doc/.ppt 老格式不支持，会返回转换指引。"
), is_safe=True)
async def doc_read(
    path: str,
    sheet: Optional[str] = None,
    page_start: Optional[int] = None,
    page_end: Optional[int] = None,
) -> str:
    """
    读取 Office/PDF 文档并提取文本。

    :param path: 文档路径（.docx/.xlsx/.xlsm/.xls/.pdf/.pptx）
    :param sheet: 工作表名称（xlsx/xls 时使用，默认第一个工作表）
    :param page_start: 起始页码（pdf 时使用，从 1 开始）
    :param page_end: 结束页码（pdf 时使用）
    """
    abs_path = os.path.abspath(path)
    ext = os.path.splitext(abs_path)[1].lower()

    if ext in _LEGACY_HINTS:
        return json.dumps(_err(
            f"不支持二进制老格式 {ext}。{_LEGACY_HINTS[ext]}",
        ), ensure_ascii=False)

    handler = _READERS.get(ext)
    if handler is None:
        return json.dumps(_err(
            f"不支持的文档格式: {ext or '(无扩展名)'}。"
            "支持 docx/xlsx/xlsm/xls/pdf/pptx；纯文本文件请用 file_read，"
            "图片请用 image_read",
            supported=sorted(_READERS),
        ), ensure_ascii=False)

    if not os.path.isfile(abs_path):
        return json.dumps(_err(f"文件不存在: {abs_path}"), ensure_ascii=False)

    try:
        result = handler(
            abs_path, sheet=sheet, page_start=page_start, page_end=page_end,
        )
        return json.dumps(result, ensure_ascii=False)
    except Exception as e:
        return json.dumps(_err(
            f"文档解析失败: {e}",
            hint="文件可能损坏、加密或格式与扩展名不符",
        ), ensure_ascii=False)
