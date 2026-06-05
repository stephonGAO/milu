"""测试内置工具 doc_read - Office/PDF 文档文本提取

- docx/xlsx/pptx 用对应库在临时目录生成真实文件后读取
- pdf 用手工构造的最小合法 PDF（自行计算 xref 偏移）
- 老格式 .doc/.ppt 返回转换指引；未知格式引导 file_read
- file_read 对二进制文档/图片扩展名返回专用工具引导
"""
import json
import tempfile
from pathlib import Path

import pytest

from milu.tools.builtin.doc_tool import (
    _MAX_SHEET_ROWS,
    doc_read,
)
from milu.tools.builtin.file_tool import file_read


@pytest.fixture
def tmp_dir():
    with tempfile.TemporaryDirectory() as d:
        yield Path(d)


def _call(**kwargs):
    return doc_read._tool_wrapper.func(**kwargs)


# ── PDF 构造辅助：最小合法单页/多页 PDF ─────────────────────


def _build_pdf_bytes(page_texts: list[str]) -> bytes:
    """构造最小合法 PDF：对象编号 1=Catalog, 2=Pages, 3..=Page/Content 交替, 末尾=Font。"""
    n = len(page_texts)
    font_num = 2 + 2 * n + 1
    objs: dict[int, str] = {
        1: "<< /Type /Catalog /Pages 2 0 R >>",
    }
    kids = []
    for i, text in enumerate(page_texts):
        page_num = 3 + i * 2
        content_num = page_num + 1
        kids.append(f"{page_num} 0 R")
        objs[page_num] = (
            f"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
            f"/Contents {content_num} 0 R "
            f"/Resources << /Font << /F1 {font_num} 0 R >> >> >>"
        )
        stream = f"BT /F1 24 Tf 72 720 Td ({text}) Tj ET"
        objs[content_num] = (
            f"<< /Length {len(stream)} >>\nstream\n{stream}\nendstream"
        )
    objs[2] = f"<< /Type /Pages /Kids [{' '.join(kids)}] /Count {n} >>"
    objs[font_num] = "<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>"

    out = b"%PDF-1.4\n"
    offsets: dict[int, int] = {}
    for num in sorted(objs):
        offsets[num] = len(out)
        out += f"{num} 0 obj\n{objs[num]}\nendobj\n".encode("latin-1")
    xref_pos = len(out)
    max_num = max(objs)
    out += f"xref\n0 {max_num + 1}\n".encode()
    out += b"0000000000 65535 f \n"
    for num in range(1, max_num + 1):
        out += f"{offsets[num]:010d} 00000 n \n".encode()
    out += (
        f"trailer\n<< /Size {max_num + 1} /Root 1 0 R >>\n"
        f"startxref\n{xref_pos}\n%%EOF\n"
    ).encode()
    return out


# ── docx ──────────────────────────────────────────────────


class TestDocx:

    async def test_paragraphs_and_tables(self, tmp_dir):
        from docx import Document
        p = tmp_dir / "t.docx"
        doc = Document()
        doc.add_paragraph("第一段内容")
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "姓名"
        table.cell(0, 1).text = "年龄"
        table.cell(1, 0).text = "张三"
        table.cell(1, 1).text = "30"
        doc.add_paragraph("结尾段落")
        doc.save(str(p))

        result = json.loads(await _call(path=str(p)))
        assert result["success"] is True
        assert result["format"] == "docx"
        assert result["paragraphs"] == 2
        assert result["tables"] == 1
        assert "第一段内容" in result["content"]
        assert "| 姓名 | 年龄 |" in result["content"]  # 表格转 Markdown
        assert "结尾段落" in result["content"]
        # 文档顺序：段落 → 表格 → 段落
        assert result["content"].index("第一段内容") \
            < result["content"].index("张三") \
            < result["content"].index("结尾段落")


# ── xlsx ──────────────────────────────────────────────────


def _make_xlsx(path: Path):
    from openpyxl import Workbook
    wb = Workbook()
    ws1 = wb.active
    ws1.title = "数据"
    ws1.append(["名称", "数量"])
    ws1.append(["苹果", 3])
    ws2 = wb.create_sheet("备注")
    ws2.append(["仅备注页"])
    wb.save(str(path))


class TestXlsx:

    async def test_default_first_sheet(self, tmp_dir):
        p = tmp_dir / "t.xlsx"
        _make_xlsx(p)
        result = json.loads(await _call(path=str(p)))
        assert result["success"] is True
        assert result["format"] == "xlsx"
        assert result["sheets"] == ["数据", "备注"]
        assert result["sheet"] == "数据"
        assert "名称\t数量" in result["content"]
        assert "苹果\t3" in result["content"]
        assert "sheet 参数" in result.get("hint", "")  # 多 sheet 提示

    async def test_named_sheet(self, tmp_dir):
        p = tmp_dir / "t.xlsx"
        _make_xlsx(p)
        result = json.loads(await _call(path=str(p), sheet="备注"))
        assert result["sheet"] == "备注"
        assert "仅备注页" in result["content"]

    async def test_sheet_not_found(self, tmp_dir):
        p = tmp_dir / "t.xlsx"
        _make_xlsx(p)
        result = json.loads(await _call(path=str(p), sheet="不存在"))
        assert result["success"] is False
        assert result["sheets"] == ["数据", "备注"]

    async def test_row_truncation(self, tmp_dir):
        from openpyxl import Workbook
        p = tmp_dir / "big.xlsx"
        wb = Workbook()
        ws = wb.active
        for i in range(_MAX_SHEET_ROWS + 10):
            ws.append([f"行{i}"])
        wb.save(str(p))

        result = json.loads(await _call(path=str(p)))
        assert result["success"] is True
        assert result["rows_read"] == _MAX_SHEET_ROWS
        assert result["truncated"] is True


# ── pdf ───────────────────────────────────────────────────


class TestPdf:

    async def test_extract_all_pages(self, tmp_dir):
        p = tmp_dir / "t.pdf"
        p.write_bytes(_build_pdf_bytes(["Hello Page One", "Second Page"]))
        result = json.loads(await _call(path=str(p)))
        assert result["success"] is True
        assert result["format"] == "pdf"
        assert result["total_pages"] == 2
        assert result["page_start"] == 1 and result["page_end"] == 2
        assert "Hello Page One" in result["content"]
        assert "Second Page" in result["content"]
        assert "--- 第 1 页 ---" in result["content"]

    async def test_page_range(self, tmp_dir):
        p = tmp_dir / "t.pdf"
        p.write_bytes(_build_pdf_bytes(["P-one", "P-two", "P-three"]))
        result = json.loads(await _call(path=str(p), page_start=2, page_end=2))
        assert result["page_start"] == 2 and result["page_end"] == 2
        assert "P-two" in result["content"]
        assert "P-one" not in result["content"]
        assert result["truncated"] is True  # 未读完全部页
        assert "page_start" in result["hint"]

    async def test_page_out_of_range(self, tmp_dir):
        p = tmp_dir / "t.pdf"
        p.write_bytes(_build_pdf_bytes(["only"]))
        result = json.loads(await _call(path=str(p), page_start=5))
        assert result["success"] is False
        assert "越界" in result["error"]


# ── pptx ──────────────────────────────────────────────────


class TestPptx:

    async def test_slides_text(self, tmp_dir):
        from pptx import Presentation
        from pptx.util import Inches
        p = tmp_dir / "t.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        box.text_frame.text = "演示文稿正文"
        prs.save(str(p))

        result = json.loads(await _call(path=str(p)))
        assert result["success"] is True
        assert result["format"] == "pptx"
        assert result["slides"] == 1
        assert "--- 幻灯片 1 ---" in result["content"]
        assert "演示文稿正文" in result["content"]


# ── 路由与错误 ────────────────────────────────────────────


class TestRouting:

    async def test_legacy_doc_hint(self, tmp_dir):
        p = tmp_dir / "old.doc"
        p.write_bytes(b"\xd0\xcf\x11\xe0")  # OLE2 头
        result = json.loads(await _call(path=str(p)))
        assert result["success"] is False
        assert "另存为 .docx" in result["error"]

    async def test_unknown_extension(self, tmp_dir):
        p = tmp_dir / "a.txt"
        p.write_text("text")
        result = json.loads(await _call(path=str(p)))
        assert result["success"] is False
        assert "file_read" in result["error"]

    async def test_file_not_exists(self):
        result = json.loads(await _call(path="Z:/no/such.docx"))
        assert result["success"] is False
        assert "不存在" in result["error"]

    async def test_corrupted_file(self, tmp_dir):
        p = tmp_dir / "bad.docx"
        p.write_bytes(b"not a zip at all")
        result = json.loads(await _call(path=str(p)))
        assert result["success"] is False
        assert "解析失败" in result["error"]


class TestFileReadRedirect:
    """file_read 对二进制格式的专用工具引导"""

    async def test_docx_redirects_to_doc_read(self, tmp_dir):
        p = tmp_dir / "t.docx"
        p.write_bytes(b"PK")
        result = json.loads(
            await file_read._tool_wrapper.func(action="read", path=str(p))
        )
        assert result["success"] is False
        assert "doc_read" in result["hint"]

    async def test_png_redirects_to_image_read(self, tmp_dir):
        p = tmp_dir / "t.png"
        p.write_bytes(b"\x89PNG")
        result = json.loads(
            await file_read._tool_wrapper.func(action="index", path=str(p))
        )
        assert result["success"] is False
        assert "image_read" in result["hint"]

    async def test_binary_unicode_error_friendly(self, tmp_dir):
        p = tmp_dir / "t.bin"
        p.write_bytes(bytes(range(256)) * 4)
        result = json.loads(
            await file_read._tool_wrapper.func(action="read", path=str(p))
        )
        assert result["success"] is False
        assert "二进制" in result["error"]
